import streamlit as st
import fitz  # PyMuPDF
from PIL import Image, ImageOps, ImageChops
import io
import os
import tempfile
import zipfile
from pdf2docx import Converter
from rembg import remove
from collections import defaultdict
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import pdfplumber

# -------------- PAGE CONFIG & CSS ----------------
st.set_page_config(page_title="PDF Studio Pro", layout="wide", page_icon="💠", initial_sidebar_state="expanded")

st.markdown("""
<style>
/* Premium iOS Light Theme */
.stApp {
    background-color: #F2F2F7;
    color: #1C1C1E;
    font-family: -apple-system, BlinkMacSystemFont, "SF Pro Display", "Segoe UI", Roboto, Helvetica, Arial, sans-serif;
}
[data-testid="stSidebar"] {
    background: rgba(255, 255, 255, 0.7) !important;
    backdrop-filter: blur(20px);
    -webkit-backdrop-filter: blur(20px);
    border-right: 1px solid rgba(0,0,0,0.1);
}
/* iOS Blue Buttons */
div.stButton > button:first-child {
    background-color: #007AFF;
    color: white;
    font-weight: 600;
    font-size: 16px;
    border: none;
    border-radius: 12px;
    padding: 0.6rem 1.4rem;
    box-shadow: 0 4px 10px rgba(0, 122, 255, 0.2);
    transition: all 0.2s ease;
}
div.stButton > button:first-child:hover {
    background-color: #0056b3;
    transform: scale(0.98);
}
/* Headings */
h1, h2, h3 {
    color: #000000;
    font-weight: 700;
    letter-spacing: -0.5px;
}
/* Cards & Uploaders */
[data-testid="stFileUploadDropzone"] {
    background-color: #FFFFFF;
    border: 2px dashed #C7C7CC;
    border-radius: 14px;
}
[data-testid="stFileUploadDropzone"]:hover {
    border-color: #007AFF;
}
/* Metrics/Info panels */
[data-testid="stMetricValue"] {
    color: #007AFF;
}
/* Custom Inputs */
.stTextInput>div>div>input, .stSelectbox>div>div>div {
    background-color: #FFFFFF;
    border-radius: 10px;
    border: 1px solid #E5E5EA;
}
</style>
""", unsafe_allow_html=True)

# -------------- HELPER FUNCTIONS ----------------
def save_uploaded_file(uploaded_file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            return tmp_file.name
    except Exception as e:
        st.error(f"Error saving file: {e}")
        return None

def display_sizes(orig_path, new_path):
    orig_size = os.path.getsize(orig_path) / (1024 * 1024)
    new_size = os.path.getsize(new_path) / (1024 * 1024)
    cols = st.columns(3)
    cols[0].metric("Original Size", f"{orig_size:.2f} MB")
    cols[1].metric("New Size", f"{new_size:.2f} MB")
    diff = orig_size - new_size
    cols[2].metric("Space Saved", f"{diff:.2f} MB", delta=f"{-diff:.2f} MB", delta_color="inverse")

# -------------- SIDEBAR NAVIGATION ----------------
st.sidebar.title("💠 PDF Studio Pro")
st.sidebar.markdown("The ultimate, all-in-one PDF toolkit.")

app_mode = st.sidebar.radio("Select Tool", [
    "🔄 Universal Format Converter",
    "💧 Watermark & Deep Clean",
    "✂️ Auto-Crop & Margin Resize",
    "📚 Merge & Split PDFs",
    "🗜️ Ultimate PDF Compressor",
    "🎨 Dark Mode Color Inverter",
    "🖼️ AI Background Remover",
    "🔐 Security (Protect/Unlock)"
])

st.sidebar.markdown("---")
st.sidebar.markdown("💡 **Tip:** All operations are performed entirely locally on your machine.")
st.sidebar.markdown("<br><br><div style='text-align: center; color: #8E8E93; font-size: 14px; font-weight: 500;'>Made with ❤️ by <b>SONU VERMA</b></div>", unsafe_allow_html=True)

# -------------- MAIN APP AREA ----------------
st.title(app_mode)

if app_mode == "🔄 Universal Format Converter":
    st.markdown("Convert between PDF and virtually any other format.")
    
    conv_type = st.selectbox("Select Conversion Type:", [
        "PDF to Word (.docx)",
        "PDF to PowerPoint (.pptx)",
        "PDF to Excel/CSV Tables (.xlsx)",
        "PDF to Raw Text (.txt)",
        "PDF to HTML (.html)",
        "PDF to High-Res Images (.zip)",
        "Images to PDF (.pdf)"
    ])
    
    if conv_type == "Images to PDF (.pdf)":
        images = st.file_uploader("Upload Images (PNG/JPG)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)
        if images and st.button("Create PDF"):
            with st.spinner("Converting images into a single PDF..."):
                try:
                    out_pdf = fitz.open()
                    for img_file in images:
                        img_bytes = img_file.getvalue()
                        # Use fitz open appropriately based on extension
                        img_doc = fitz.open("jpeg", img_bytes) if img_file.name.lower().endswith(('jpg','jpeg')) else fitz.open("png", img_bytes)
                        pdf_bytes = img_doc.convert_to_pdf()
                        out_pdf.insert_pdf(fitz.open("pdf", pdf_bytes))
                    
                    out_path = "converted_images.pdf"
                    out_pdf.save(out_path, deflate=True)
                    out_pdf.close()
                    with open(out_path, "rb") as f:
                        st.download_button("Download PDF", f, file_name="images.pdf")
                except Exception as e:
                    st.error(f"Error: {e}")
    else:
        conv_file = st.file_uploader("Upload PDF to convert", type=["pdf"])
        if conv_file and st.button("Start Conversion"):
            with st.spinner(f"Running conversion algorithms..."):
                path = save_uploaded_file(conv_file)
                base_name = os.path.splitext(conv_file.name)[0]
                
                try:
                    if "Word" in conv_type:
                        out = path.replace(".pdf", ".docx")
                        cv = Converter(path)
                        cv.convert(out, start=0, end=None)
                        cv.close()
                        with open(out, "rb") as f:
                            st.download_button("Download Word Doc", f, file_name=f"{base_name}.docx")
                            
                    elif "PowerPoint" in conv_type:
                        out = path.replace(".pdf", ".pptx")
                        doc = fitz.open(path)
                        prs = Presentation()
                        # Standard 16:9 Slide
                        prs.slide_width = Inches(16)
                        prs.slide_height = Inches(9)
                        blank_slide_layout = prs.slide_layouts[6] # 6 is usually a blank layout
                        
                        for page in doc:
                            pix = page.get_pixmap(dpi=150)
                            img_data = pix.tobytes("jpeg")
                            img_stream = io.BytesIO(img_data)
                            slide = prs.slides.add_slide(blank_slide_layout)
                            # Stretch image to fit slide fully
                            slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                            
                        prs.save(out)
                        with open(out, "rb") as f:
                            st.download_button("Download PowerPoint", f, file_name=f"{base_name}.pptx")
                            
                    elif "Excel" in conv_type:
                        out = path.replace(".pdf", ".xlsx")
                        with pdfplumber.open(path) as pdf:
                            all_tables = []
                            for i, page in enumerate(pdf.pages):
                                tables = page.extract_tables()
                                for table in tables:
                                    if table and len(table) > 1:
                                        # Use first row as header if it exists
                                        header = table[0]
                                        # Replace None with empty string in header to avoid pandas errors
                                        header = [str(x) if x is not None else f"Col_{j}" for j, x in enumerate(header)]
                                        df = pd.DataFrame(table[1:], columns=header)
                                        all_tables.append(df)
                                        
                        if all_tables:
                            with pd.ExcelWriter(out) as writer:
                                for i, df in enumerate(all_tables):
                                    df.to_excel(writer, sheet_name=f"Table_{i+1}", index=False)
                            with open(out, "rb") as f:
                                st.download_button("Download Excel", f, file_name=f"{base_name}.xlsx")
                        else:
                            st.warning("No structured tables were detected in the PDF!")
                            
                    elif "Text" in conv_type:
                        out = path.replace(".pdf", ".txt")
                        doc = fitz.open(path)
                        text = ""
                        for i, page in enumerate(doc):
                            text += f"--- Page {i+1} ---\n\n"
                            text += page.get_text("text") + "\n\n"
                        with open(out, "w", encoding="utf-8") as f:
                            f.write(text)
                        with open(out, "rb") as f:
                            st.download_button("Download Text File", f, file_name=f"{base_name}.txt")
                            
                    elif "HTML" in conv_type:
                        out = path.replace(".pdf", ".html")
                        doc = fitz.open(path)
                        html = "<html><head><meta charset='utf-8'></head><body>\n"
                        for page in doc:
                            html += page.get_text("html") + "<hr/>\n"
                        html += "</body></html>"
                        with open(out, "w", encoding="utf-8") as f:
                            f.write(html)
                        with open(out, "rb") as f:
                            st.download_button("Download HTML", f, file_name=f"{base_name}.html")
                            
                    elif "Images" in conv_type:
                        out = path.replace(".pdf", ".zip")
                        doc = fitz.open(path)
                        with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zipf:
                            for i, page in enumerate(doc):
                                pix = page.get_pixmap(dpi=300)
                                zipf.writestr(f"page_{i+1}.jpg", pix.tobytes("jpeg"))
                        with open(out, "rb") as f:
                            st.download_button("Download Images (ZIP)", f, file_name=f"{base_name}_images.zip")
                            
                except Exception as e:
                    st.error(f"Error during conversion: {e}")

elif app_mode == "💧 Watermark & Deep Clean":
    st.markdown("Advanced multi-layered watermark extraction. Use **Deep Clean** if auto-detect fails.")
    wm_file = st.file_uploader("Upload PDF", type=["pdf"])
    
    mode = st.selectbox("Select Cleaning Engine:", [
        "1. Remove Annotations & Stamps (Best for invisible/stubborn watermarks)",
        "2. Auto-Detect Repeating Text",
        "3. Manual Text Target",
        "4. Remove All Background Images",
        "5. Extreme: Remove ALL Text (Keep only graphics)"
    ])
    
    wm_text = ""
    if mode == "3. Manual Text Target":
        wm_text = st.text_input("Exact Text to Remove (Case-sensitive):")
        
    if wm_file and st.button("Clean Document"):
        with st.spinner("Executing Deep Clean algorithms..."):
            input_path = save_uploaded_file(wm_file)
            output_path = input_path.replace(".pdf", "_clean.pdf")
            
            try:
                doc = fitz.open(input_path)
                count = 0
                
                if mode == "1. Remove Annotations & Stamps (Best for invisible/stubborn watermarks)":
                    for page in doc:
                        for annot in page.annots():
                            page.delete_annot(annot)
                            count += 1
                            
                elif mode == "3. Manual Text Target" and wm_text:
                    for page in doc:
                        for inst in page.search_for(wm_text):
                            page.add_redact_annot(inst, fill=(1, 1, 1))
                            count += 1
                        page.apply_redactions()
                        
                elif mode == "2. Auto-Detect Repeating Text":
                    text_counts = defaultdict(int)
                    num_pages = len(doc)
                    if num_pages > 1:
                        for page in doc:
                            blocks = page.get_text("blocks")
                            seen_on_page = set()
                            for b in blocks:
                                text = b[4].strip()
                                if len(text) > 3 and text not in seen_on_page:
                                    text_counts[text] += 1
                                    seen_on_page.add(text)
                        
                        candidates = [t for t, c in text_counts.items() if c >= (num_pages * 0.5)]
                        for page in doc:
                            for candidate in candidates:
                                for inst in page.search_for(candidate):
                                    page.add_redact_annot(inst, fill=(1, 1, 1))
                                    count += 1
                            page.apply_redactions()
                        if count > 0:
                            st.success("Auto-removed repeating text watermarks!")
                        else:
                            st.warning("No repeating text found. Try removing Annotations instead.")
                    else:
                        st.warning("Needs >1 page for auto-detect. Use Manual mode.")

                elif mode == "4. Remove All Background Images":
                    for page in doc:
                        for img in page.get_images(full=True):
                            page.delete_image(img[0])
                            count += 1
                            
                elif mode == "5. Extreme: Remove ALL Text (Keep only graphics)":
                    for page in doc:
                        for b in page.get_text("blocks"):
                            if b[6] == 0:
                                page.add_redact_annot(fitz.Rect(b[:4]))
                                count += 1
                        page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
                        
                doc.save(output_path, deflate=True, garbage=4)
                doc.close()
                
                st.success(f"Processing Complete! Removed {count} artifacts.")
                display_sizes(input_path, output_path)
                with open(output_path, "rb") as f:
                    st.download_button("Download Cleaned PDF", f, file_name=wm_file.name.replace(".pdf", "_cleaned.pdf"), mime="application/pdf")
            except Exception as e:
                st.error(f"Error: {e}")

elif app_mode == "✂️ Auto-Crop & Margin Resize":
    st.markdown("Automatically detect and slice away blank white boundaries around your document, or set manual margins.")
    crop_file = st.file_uploader("Upload PDF", type=["pdf"])
    
    crop_mode = st.radio("Cropping Method", ["Auto-Detect Blank Boundaries", "Manual Margins (Points)"])
    
    if crop_mode == "Manual Margins (Points)":
        cols = st.columns(4)
        c_top = cols[0].number_input("Top Margin", value=50)
        c_bot = cols[1].number_input("Bottom Margin", value=50)
        c_left = cols[2].number_input("Left Margin", value=50)
        c_right = cols[3].number_input("Right Margin", value=50)
        
    if crop_file and st.button("Crop PDF"):
        with st.spinner("Analyzing margins and cropping..."):
            input_path = save_uploaded_file(crop_file)
            output_path = input_path.replace(".pdf", "_cropped.pdf")
            
            try:
                doc = fitz.open(input_path)
                
                if crop_mode == "Auto-Detect Blank Boundaries":
                    progress = st.progress(0)
                    for i, page in enumerate(doc):
                        pix = page.get_pixmap(dpi=72)
                        img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                        gray = img.convert("L")
                        inverted = ImageOps.invert(gray)
                        bbox = inverted.getbbox() 
                        
                        if bbox:
                            pad = 10
                            rect = fitz.Rect(bbox[0]-pad, bbox[1]-pad, bbox[2]+pad, bbox[3]+pad)
                            mb = page.rect
                            rect.intersect(mb) 
                            page.set_cropbox(rect)
                        progress.progress((i+1)/len(doc))
                else:
                    for page in doc:
                        rect = page.rect
                        new_rect = fitz.Rect(rect.x0 + c_left, rect.y0 + c_top, rect.x1 - c_right, rect.y1 - c_bot)
                        page.set_cropbox(new_rect)
                        
                doc.save(output_path, deflate=True, garbage=4)
                doc.close()
                st.success("Successfully cropped all pages!")
                with open(output_path, "rb") as f:
                    st.download_button("Download Cropped PDF", f, file_name="cropped.pdf")
            except Exception as e:
                st.error(f"Error: {e}")

elif app_mode == "📚 Merge & Split PDFs":
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("Merge PDFs")
        merge_files = st.file_uploader("Upload multiple PDFs to combine", type=["pdf"], accept_multiple_files=True)
        if merge_files and st.button("Merge Files"):
            with st.spinner("Merging..."):
                try:
                    result = fitz.open()
                    for f in merge_files:
                        path = save_uploaded_file(f)
                        doc = fitz.open(path)
                        result.insert_pdf(doc)
                        doc.close()
                    out_path = "merged_output.pdf"
                    result.save(out_path, deflate=True, garbage=4)
                    result.close()
                    st.success("Merged successfully!")
                    with open(out_path, "rb") as f:
                        st.download_button("Download Merged PDF", f, file_name="merged.pdf")
                except Exception as e:
                    st.error(f"Error: {e}")

    with col2:
        st.subheader("Split PDF")
        split_file = st.file_uploader("Upload PDF to Split", type=["pdf"], key="split")
        pages_input = st.text_input("Pages to extract (e.g. '1-3, 5, 8-10'):")
        if split_file and pages_input and st.button("Extract Pages"):
            with st.spinner("Splitting..."):
                try:
                    path = save_uploaded_file(split_file)
                    doc = fitz.open(path)
                    result = fitz.open()
                    
                    pages_to_keep = []
                    parts = pages_input.split(',')
                    for p in parts:
                        p = p.strip()
                        if '-' in p:
                            start, end = map(int, p.split('-'))
                            pages_to_keep.extend(list(range(start-1, end)))
                        else:
                            pages_to_keep.append(int(p)-1)
                            
                    for p_num in pages_to_keep:
                        if 0 <= p_num < len(doc):
                            result.insert_pdf(doc, from_page=p_num, to_page=p_num)
                            
                    out_path = "split_output.pdf"
                    result.save(out_path, deflate=True, garbage=4)
                    doc.close()
                    result.close()
                    st.success("Extracted successfully!")
                    with open(out_path, "rb") as f:
                        st.download_button("Download Extracted Pages", f, file_name="extracted.pdf")
                except Exception as e:
                    st.error(f"Error: {e}")

elif app_mode == "🗜️ Ultimate PDF Compressor":
    st.markdown("Massively reduce PDF file size by destroying unused objects, flattening streams, and downscaling images.")
    comp_file = st.file_uploader("Upload PDF to compress", type=["pdf"])
    
    downscale = st.checkbox("Aggressive: Downscale all embedded images to 72 DPI (Drastic size reduction, lower image quality)", value=False)
    
    if comp_file and st.button("Compress Document"):
        with st.spinner("Compressing..."):
            input_path = save_uploaded_file(comp_file)
            output_path = input_path.replace(".pdf", "_compressed.pdf")
            
            try:
                doc = fitz.open(input_path)
                
                if downscale:
                    for page in doc:
                        for img in page.get_images(full=True):
                            xref = img[0]
                            base_img = doc.extract_image(xref)
                            img_data = base_img["image"]
                            
                            pil_img = Image.open(io.BytesIO(img_data))
                            if pil_img.mode in ("RGBA", "P"):
                                pil_img = pil_img.convert("RGB")
                                
                            w, h = pil_img.size
                            pil_img = pil_img.resize((int(w*0.5), int(h*0.5)), Image.Resampling.LANCZOS)
                            
                            b = io.BytesIO()
                            pil_img.save(b, format="JPEG", quality=60, optimize=True)
                            page.insert_image(page.rect, stream=b.getvalue(), replace=True)
                            
                doc.save(output_path, garbage=4, deflate=True, clean=True)
                doc.close()
                st.success("Compression Complete!")
                display_sizes(input_path, output_path)
                with open(output_path, "rb") as f:
                    st.download_button("Download Compressed PDF", f, file_name="compressed.pdf")
            except Exception as e:
                st.error(f"Error: {e}")

elif app_mode == "🎨 Dark Mode Color Inverter":
    st.markdown("Converts all pages to inverted dark-mode images.")
    inv_file = st.file_uploader("Upload PDF", type=["pdf"])
    dpi_setting = st.slider("Render Quality (DPI)", 72, 300, 150)
    if inv_file and st.button("Invert PDF"):
        with st.spinner("Inverting..."):
            input_path = save_uploaded_file(inv_file)
            out_path = input_path.replace(".pdf", "_dark.pdf")
            try:
                doc = fitz.open(input_path)
                out_pdf = fitz.open()
                bar = st.progress(0)
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=dpi_setting)
                    img = ImageOps.invert(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
                    b = io.BytesIO()
                    img.save(b, format="JPEG", quality=80)
                    out_pdf.insert_pdf(fitz.open("pdf", fitz.open("jpeg", b.getvalue()).convert_to_pdf()))
                    bar.progress((i+1)/len(doc))
                out_pdf.save(out_path, deflate=True, garbage=4)
                st.success("Dark Mode ready!")
                display_sizes(input_path, out_path)
                with open(out_path, "rb") as f:
                    st.download_button("Download", f, file_name="dark_mode.pdf")
            except Exception as e:
                st.error(f"Error: {e}")

elif app_mode == "🖼️ AI Background Remover":
    st.markdown("Isolates foreground elements using U2Net Deep Learning.")
    bg_file = st.file_uploader("Upload PDF", type=["pdf"])
    if bg_file and st.button("Run AI Separation"):
        with st.spinner("Running AI Model..."):
            input_path = save_uploaded_file(bg_file)
            out_path = input_path.replace(".pdf", "_nobg.pdf")
            try:
                doc = fitz.open(input_path)
                out_pdf = fitz.open()
                bar = st.progress(0)
                for i, page in enumerate(doc):
                    pix = page.get_pixmap(dpi=150)
                    img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    no_bg = remove(img)
                    bg = Image.new("RGB", no_bg.size, (255, 255, 255))
                    bg.paste(no_bg, mask=no_bg.split()[3] if len(no_bg.split())==4 else None)
                    b = io.BytesIO()
                    bg.save(b, format="JPEG", quality=80)
                    out_pdf.insert_pdf(fitz.open("pdf", fitz.open("jpeg", b.getvalue()).convert_to_pdf()))
                    bar.progress((i+1)/len(doc))
                out_pdf.save(out_path, deflate=True, garbage=4)
                st.success("Backgrounds eradicated!")
                with open(out_path, "rb") as f:
                    st.download_button("Download", f, file_name="no_background.pdf")
            except Exception as e:
                st.error(f"Error: {e}")

elif app_mode == "🔐 Security (Protect/Unlock)":
    st.markdown("Encrypt your PDF with a password, or remove password protection.")
    sec_file = st.file_uploader("Upload PDF", type=["pdf"])
    action = st.radio("Action", ["Add Password", "Remove Password"])
    pwd = st.text_input("Password:", type="password")
    
    if sec_file and pwd and st.button("Execute"):
        with st.spinner("Processing security protocol..."):
            input_path = save_uploaded_file(sec_file)
            out_path = "secured.pdf"
            try:
                doc = fitz.open(input_path)
                if action == "Remove Password":
                    if doc.authenticate(pwd):
                        doc.save(out_path)
                        st.success("Unlocked successfully!")
                    else:
                        st.error("Incorrect Password!")
                        st.stop()
                else:
                    doc.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=pwd, owner_pw=pwd)
                    st.success("Encrypted successfully!")
                doc.close()
                with open(out_path, "rb") as f:
                    st.download_button("Download File", f, file_name="secured_doc.pdf")
            except Exception as e:
                st.error(f"Error: {e}")
