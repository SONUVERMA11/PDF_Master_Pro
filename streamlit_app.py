import streamlit as st
import fitz  # PyMuPDF
from PIL import Image, ImageOps, ImageChops
import io
import os
import tempfile
import zipfile
from pdf2docx import Converter
from rembg import remove
from collections import defaultdict
import pandas as pd
from pptx import Presentation
from pptx.util import Inches
import pdfplumber

# -------------- PAGE CONFIG & CSS ----------------
st.set_page_config(page_title="PDF Studio Pro", layout="wide", page_icon="💠", initial_sidebar_state="collapsed")

st.markdown("""
<style>
/* ----------------------------------------------------
   HIGH CONTRAST DARK THEME - NEON/CYBERPUNK
   ---------------------------------------------------- */
.stApp {
    background-color: #050505;
    color: #e0e0e0;
    font-family: "SF Pro Display", "Inter", sans-serif;
}

/* Hide Sidebar since we use Tabs now */
[data-testid="collapsedControl"] { display: none; }
section[data-testid="stSidebar"] { display: none; }

/* Tabs Styling (Streamlit native tabs) */
div[data-testid="stTabs"] {
    background-color: #050505;
}
div[data-testid="stTabs"] button {
    color: #888888 !important;
    font-weight: 600;
    font-size: 1.1rem;
    padding: 12px 24px;
    background: transparent !important;
    border: none !important;
    border-bottom: 2px solid transparent !important;
}
div[data-testid="stTabs"] button[aria-selected="true"] {
    color: #00f0ff !important;
    border-bottom: 2px solid #00f0ff !important;
    text-shadow: 0 0 10px rgba(0, 240, 255, 0.6);
}

/* Stunning Neon Buttons */
div.stButton > button:first-child {
    background: #0a0a0a;
    color: #00f0ff;
    font-weight: 800;
    font-size: 1rem;
    text-transform: uppercase;
    letter-spacing: 1.5px;
    border: 1px solid #00f0ff;
    border-radius: 4px;
    padding: 0.8rem 2rem;
    box-shadow: 0 0 8px rgba(0, 240, 255, 0.2), inset 0 0 8px rgba(0, 240, 255, 0.1);
    transition: all 0.2s ease-in-out;
    width: 100%;
}
div.stButton > button:first-child:hover {
    background: #00f0ff;
    color: #000000;
    box-shadow: 0 0 20px rgba(0, 240, 255, 0.8), inset 0 0 15px rgba(0, 240, 255, 0.6);
    transform: scale(1.02);
}

/* Headings */
h1, h2, h3 {
    color: #ffffff;
    font-weight: 900;
    text-transform: uppercase;
    letter-spacing: 2px;
}
h1 { 
    font-size: 3rem; 
    margin-bottom: 0.5rem; 
    text-shadow: 3px 3px 0px #ff003c;
}
h2 {
    color: #00f0ff;
    border-bottom: 1px solid #222;
    padding-bottom: 10px;
}

/* File Uploader */
[data-testid="stFileUploadDropzone"] {
    background-color: #0a0a0a;
    border: 2px dashed #333333;
    border-radius: 8px;
    padding: 40px;
    transition: all 0.3s ease;
}
[data-testid="stFileUploadDropzone"]:hover {
    border-color: #ff003c;
    background-color: rgba(255, 0, 60, 0.05);
    box-shadow: 0 0 20px rgba(255, 0, 60, 0.2);
}

/* Inputs & Selectors */
.stTextInput>div>div>input, .stSelectbox>div>div>div, .stNumberInput>div>div>input {
    background-color: #0a0a0a;
    color: #ffffff;
    border-radius: 4px;
    border: 1px solid #333;
    padding: 12px 16px;
    font-size: 15px;
    transition: all 0.3s ease;
}
.stTextInput>div>div>input:focus, .stSelectbox>div>div>div:focus, .stNumberInput>div>div>input:focus {
    border-color: #00f0ff;
    box-shadow: 0 0 10px rgba(0, 240, 255, 0.3);
    color: #00f0ff;
}

/* Metrics & Stats */
[data-testid="stMetricValue"] {
    font-size: 2.5rem;
    font-weight: 800;
    color: #00f0ff;
    text-shadow: 0 0 10px rgba(0,240,255,0.4);
}
[data-testid="stMetricDelta"] {
    color: #ff003c !important;
}

/* Container Cards */
[data-testid="stVerticalBlock"] > [style*="flex-direction: column;"] > [data-testid="stVerticalBlock"] {
    background: #0a0a0a;
    border-radius: 8px;
    padding: 30px;
    border: 1px solid #1c1c1c;
    box-shadow: 0 10px 30px rgba(0,0,0,0.8);
    animation: slideUp 0.4s ease-out forwards;
}

@keyframes slideUp {
    from { opacity: 0; transform: translateY(10px); }
    to { opacity: 1; transform: translateY(0); }
}

/* Footer */
.custom-footer {
    position: fixed;
    bottom: 0; left: 0; width: 100%;
    text-align: center;
    padding: 15px;
    background: rgba(5, 5, 5, 0.9);
    backdrop-filter: blur(10px);
    color: #888;
    border-top: 1px solid #222;
    font-size: 15px;
    z-index: 9999;
    letter-spacing: 1px;
}
.custom-footer span {
    color: #ff003c;
    font-weight: 900;
    text-shadow: 0 0 8px rgba(255,0,60,0.6);
}

/* Hide clutter */
#MainMenu {visibility: hidden;}
footer {visibility: hidden;}
header {visibility: hidden;}
</style>
""", unsafe_allow_html=True)

# -------------- HELPER FUNCTIONS ----------------
def save_uploaded_file(uploaded_file):
    try:
        with tempfile.NamedTemporaryFile(delete=False, suffix=".pdf") as tmp_file:
            tmp_file.write(uploaded_file.getvalue())
            return tmp_file.name
    except Exception as e:
        st.error(f"Error saving file: {e}")
        return None

def display_sizes(orig_path, new_path):
    orig_size = os.path.getsize(orig_path) / (1024 * 1024)
    new_size = os.path.getsize(new_path) / (1024 * 1024)
    cols = st.columns(3)
    cols[0].metric("Original Size", f"{orig_size:.2f} MB")
    cols[1].metric("New Size", f"{new_size:.2f} MB")
    diff = orig_size - new_size
    cols[2].metric("Space Saved", f"{diff:.2f} MB", delta=f"{-diff:.2f} MB", delta_color="inverse")

# -------------- MAIN APP LAYOUT ----------------
st.markdown("<h1>PDF STUDIO PRO <span>// TERMINAL</span></h1>", unsafe_allow_html=True)
st.markdown("---")

tabs = st.tabs([
    "🔄 Converter",
    "💧 Deep Clean",
    "✂️ Crop Engine",
    "📚 Merge/Split",
    "🗜️ Compressor",
    "🎨 Inverter",
    "🖼️ Rembg AI",
    "🔐 Security"
])

# ----------------- 1. CONVERTER -----------------
with tabs[0]:
    st.header("🔄 Universal Converter")
    st.markdown("Convert between PDF and virtually any other format seamlessly.")
    
    col1, col2 = st.columns(2)
    with col1:
        conv_type = st.selectbox("Select Conversion Matrix:", [
            "PDF to Word (.docx)",
            "PDF to PowerPoint (.pptx)",
            "PDF to Excel/CSV Tables (.xlsx)",
            "PDF to Raw Text (.txt)",
            "PDF to HTML (.html)",
            "PDF to High-Res Images (.zip)",
            "Images to PDF (.pdf)"
        ])
        
        if conv_type == "Images to PDF (.pdf)":
            images = st.file_uploader("Upload Image Pack (PNG/JPG)", type=["png", "jpg", "jpeg"], accept_multiple_files=True)
            if images and st.button("EXECUTE BUILD"):
                with st.spinner("Compiling images into PDF architecture..."):
                    try:
                        out_pdf = fitz.open()
                        for img_file in images:
                            img_bytes = img_file.getvalue()
                            img_doc = fitz.open("jpeg", img_bytes) if img_file.name.lower().endswith(('jpg','jpeg')) else fitz.open("png", img_bytes)
                            out_pdf.insert_pdf(fitz.open("pdf", img_doc.convert_to_pdf()))
                        
                        out_path = "converted_images.pdf"
                        out_pdf.save(out_path, deflate=True)
                        out_pdf.close()
                        with col2:
                            st.success("COMPILE COMPLETE")
                            with open(out_path, "rb") as f:
                                st.download_button("DOWNLOAD ARTIFACT", f, file_name="images.pdf")
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")
        else:
            conv_file = st.file_uploader("Upload Source PDF", type=["pdf"])
            if conv_file and st.button("INITIALIZE CONVERSION"):
                with col2:
                    with st.spinner(f"Running conversion algorithms for {conv_type}..."):
                        path = save_uploaded_file(conv_file)
                        base_name = os.path.splitext(conv_file.name)[0]
                        try:
                            if "Word" in conv_type:
                                out = path.replace(".pdf", ".docx")
                                cv = Converter(path)
                                cv.convert(out, start=0, end=None)
                                cv.close()
                                st.success("CONVERSION COMPLETE")
                                with open(out, "rb") as f:
                                    st.download_button("DOWNLOAD WORD DOC", f, file_name=f"{base_name}.docx")
                                    
                            elif "PowerPoint" in conv_type:
                                out = path.replace(".pdf", ".pptx")
                                doc = fitz.open(path)
                                prs = Presentation()
                                prs.slide_width, prs.slide_height = Inches(16), Inches(9)
                                blank = prs.slide_layouts[6] 
                                
                                bar = st.progress(0)
                                for i, page in enumerate(doc):
                                    img_stream = io.BytesIO(page.get_pixmap(dpi=150).tobytes("jpeg"))
                                    slide = prs.slides.add_slide(blank)
                                    slide.shapes.add_picture(img_stream, 0, 0, width=prs.slide_width, height=prs.slide_height)
                                    bar.progress((i+1)/len(doc))
                                    
                                prs.save(out)
                                st.success("CONVERSION COMPLETE")
                                with open(out, "rb") as f:
                                    st.download_button("DOWNLOAD PPTX", f, file_name=f"{base_name}.pptx")
                                    
                            elif "Excel" in conv_type:
                                out = path.replace(".pdf", ".xlsx")
                                with pdfplumber.open(path) as pdf:
                                    all_tables = []
                                    for page in pdf.pages:
                                        for table in page.extract_tables():
                                            if table and len(table) > 1:
                                                header = [str(x) if x is not None else f"Col_{j}" for j, x in enumerate(table[0])]
                                                all_tables.append(pd.DataFrame(table[1:], columns=header))
                                                
                                if all_tables:
                                    with pd.ExcelWriter(out) as writer:
                                        for i, df in enumerate(all_tables):
                                            df.to_excel(writer, sheet_name=f"Data_{i+1}", index=False)
                                    st.success("DATA EXTRACTION COMPLETE")
                                    with open(out, "rb") as f:
                                        st.download_button("DOWNLOAD EXCEL", f, file_name=f"{base_name}.xlsx")
                                else:
                                    st.warning("NO TABLES DETECTED IN STREAM")
                                    
                            elif "Text" in conv_type:
                                out = path.replace(".pdf", ".txt")
                                doc = fitz.open(path)
                                text = "".join([f"--- Page {i+1} ---\n\n{page.get_text('text')}\n\n" for i, page in enumerate(doc)])
                                with open(out, "w", encoding="utf-8") as f: f.write(text)
                                st.success("EXTRACTION COMPLETE")
                                with open(out, "rb") as f:
                                    st.download_button("DOWNLOAD RAW TEXT", f, file_name=f"{base_name}.txt")
                                    
                            elif "HTML" in conv_type:
                                out = path.replace(".pdf", ".html")
                                doc = fitz.open(path)
                                html = "<html><head><meta charset='utf-8'></head><body style='background:#050505; color:#fff;'>\n"
                                html += "".join([page.get_text("html") + "<hr/>\n" for page in doc])
                                html += "</body></html>"
                                with open(out, "w", encoding="utf-8") as f: f.write(html)
                                st.success("HTML COMPILE COMPLETE")
                                with open(out, "rb") as f:
                                    st.download_button("DOWNLOAD HTML", f, file_name=f"{base_name}.html")
                                    
                            elif "Images" in conv_type:
                                out = path.replace(".pdf", ".zip")
                                doc = fitz.open(path)
                                bar = st.progress(0)
                                with zipfile.ZipFile(out, 'w', zipfile.ZIP_DEFLATED) as zipf:
                                    for i, page in enumerate(doc):
                                        zipf.writestr(f"page_{i+1}.jpg", page.get_pixmap(dpi=300).tobytes("jpeg"))
                                        bar.progress((i+1)/len(doc))
                                st.success("EXTRACTION COMPLETE")
                                with open(out, "rb") as f:
                                    st.download_button("DOWNLOAD IMAGE PACK", f, file_name=f"{base_name}_images.zip")
                        except Exception as e:
                            st.error(f"SYSTEM FAULT: {e}")

# ----------------- 2. DEEP CLEAN / WATERMARK -----------------
with tabs[1]:
    st.header("💧 Deep Clean Engine")
    
    col1, col2 = st.columns(2)
    with col1:
        wm_file = st.file_uploader("Upload PDF Target", type=["pdf"])
        mode = st.selectbox("Select Purge Protocol:", [
            "Remove Annotations & Stamps",
            "Auto-Detect Repeating Text",
            "Manual Text Target",
            "Purge All Background Images",
            "Extreme: Purge ALL Text"
        ])
        
        wm_text = ""
        if mode == "Manual Text Target":
            wm_text = st.text_input("Target String (Case-sensitive):")
            
        if wm_file and st.button("EXECUTE PURGE"):
            with col2:
                with st.spinner("Executing Deep Clean algorithms..."):
                    input_path = save_uploaded_file(wm_file)
                    output_path = input_path.replace(".pdf", "_clean.pdf")
                    try:
                        doc = fitz.open(input_path)
                        count = 0
                        
                        if mode == "Remove Annotations & Stamps":
                            for page in doc:
                                for annot in page.annots():
                                    page.delete_annot(annot)
                                    count += 1
                                    
                        elif mode == "Manual Text Target" and wm_text:
                            for page in doc:
                                for inst in page.search_for(wm_text):
                                    page.add_redact_annot(inst, fill=None)
                                    count += 1
                                page.apply_redactions()
                                
                        elif mode == "Auto-Detect Repeating Text":
                            text_counts = defaultdict(int)
                            num_pages = len(doc)
                            if num_pages > 1:
                                for page in doc:
                                    seen = set()
                                    for b in page.get_text("blocks"):
                                        t = b[4].strip()
                                        if len(t) > 3 and t not in seen:
                                            text_counts[t] += 1
                                            seen.add(t)
                                candidates = [t for t, c in text_counts.items() if c >= (num_pages * 0.5)]
                                for page in doc:
                                    for candidate in candidates:
                                        for inst in page.search_for(candidate):
                                            page.add_redact_annot(inst, fill=None)
                                            count += 1
                                    page.apply_redactions()
                            else:
                                st.warning("REQUIRE >1 PAGE FOR AUTO-DETECT")
        
                        elif mode == "Purge All Background Images":
                            for page in doc:
                                for img in page.get_images(full=True):
                                    page.delete_image(img[0])
                                    count += 1
                                    
                        elif mode == "Extreme: Purge ALL Text":
                            for page in doc:
                                for b in page.get_text("blocks"):
                                    if b[6] == 0:
                                        page.add_redact_annot(fitz.Rect(b[:4]))
                                        count += 1
                                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
                                
                        doc.save(output_path, deflate=True, garbage=4)
                        doc.close()
                        
                        st.success(f"PURGE COMPLETE: {count} artifacts destroyed.")
                        display_sizes(input_path, output_path)
                        with open(output_path, "rb") as f:
                            st.download_button("DOWNLOAD PURGED ARTIFACT", f, file_name=wm_file.name.replace(".pdf", "_cleaned.pdf"))
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")

# ----------------- 3. CROP ENGINE -----------------
with tabs[2]:
    st.header("✂️ Auto-Crop Engine")
    col1, col2 = st.columns(2)
    
    with col1:
        crop_file = st.file_uploader("Upload Target PDF", type=["pdf"])
        crop_mode = st.radio("Targeting Method", ["Auto-Detect Empty Space", "Manual Matrix Margin"])
        
        c_top, c_bot, c_left, c_right = 50, 50, 50, 50
        if crop_mode == "Manual Matrix Margin":
            cc1, cc2 = st.columns(2)
            c_top = cc1.number_input("Top (pt)", value=50)
            c_bot = cc2.number_input("Bottom (pt)", value=50)
            c_left = cc1.number_input("Left (pt)", value=50)
            c_right = cc2.number_input("Right (pt)", value=50)
            
        if crop_file and st.button("EXECUTE CROP"):
            with col2:
                with st.spinner("Analyzing structural boundaries..."):
                    input_path = save_uploaded_file(crop_file)
                    output_path = input_path.replace(".pdf", "_cropped.pdf")
                    try:
                        doc = fitz.open(input_path)
                        if crop_mode == "Auto-Detect Empty Space":
                            bar = st.progress(0)
                            for i, page in enumerate(doc):
                                pix = page.get_pixmap(dpi=72)
                                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                                inverted = ImageOps.invert(img.convert("L"))
                                bbox = inverted.getbbox() 
                                if bbox:
                                    pad = 10
                                    rect = fitz.Rect(bbox[0]-pad, bbox[1]-pad, bbox[2]+pad, bbox[3]+pad)
                                    rect.intersect(page.rect) 
                                    page.set_cropbox(rect)
                                bar.progress((i+1)/len(doc))
                        else:
                            for page in doc:
                                r = page.rect
                                page.set_cropbox(fitz.Rect(r.x0 + c_left, r.y0 + c_top, r.x1 - c_right, r.y1 - c_bot))
                                
                        doc.save(output_path, deflate=True, garbage=4)
                        doc.close()
                        st.success("CROP EXECUTED SUCCESSFULLY")
                        with open(output_path, "rb") as f:
                            st.download_button("DOWNLOAD CROPPED PDF", f, file_name="cropped.pdf")
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")

# ----------------- 4. MERGE & SPLIT -----------------
with tabs[3]:
    st.header("📚 Matrix Operations (Merge/Split)")
    col1, col2 = st.columns(2)
    
    with col1:
        st.subheader("COMBINE")
        merge_files = st.file_uploader("Upload Segments to Merge", type=["pdf"], accept_multiple_files=True)
        if merge_files and st.button("EXECUTE MERGE"):
            with st.spinner("Compiling segments..."):
                try:
                    result = fitz.open()
                    for f in merge_files:
                        result.insert_pdf(fitz.open(save_uploaded_file(f)))
                    out_path = "merged.pdf"
                    result.save(out_path, deflate=True, garbage=4)
                    result.close()
                    st.success("MERGE SUCCESSFUL")
                    with open(out_path, "rb") as f:
                        st.download_button("DOWNLOAD COMPILED PDF", f, file_name="merged.pdf")
                except Exception as e:
                    st.error(f"SYSTEM FAULT: {e}")

    with col2:
        st.subheader("FRACTURE")
        split_file = st.file_uploader("Upload Core PDF to Split", type=["pdf"], key="split")
        pages_input = st.text_input("Target pages (e.g. 1-3, 5):")
        if split_file and pages_input and st.button("EXECUTE FRACTURE"):
            with st.spinner("Fracturing core document..."):
                try:
                    doc = fitz.open(save_uploaded_file(split_file))
                    result = fitz.open()
                    pages_to_keep = []
                    for p in pages_input.split(','):
                        p = p.strip()
                        if '-' in p:
                            start, end = map(int, p.split('-'))
                            pages_to_keep.extend(list(range(start-1, end)))
                        else:
                            pages_to_keep.append(int(p)-1)
                            
                    for p_num in pages_to_keep:
                        if 0 <= p_num < len(doc):
                            result.insert_pdf(doc, from_page=p_num, to_page=p_num)
                            
                    out_path = "split.pdf"
                    result.save(out_path, deflate=True, garbage=4)
                    doc.close()
                    st.success("FRACTURE SUCCESSFUL")
                    with open(out_path, "rb") as f:
                        st.download_button("DOWNLOAD FRAGMENT", f, file_name="extracted.pdf")
                except Exception as e:
                    st.error(f"SYSTEM FAULT: {e}")

# ----------------- 5. COMPRESS -----------------
with tabs[4]:
    st.header("🗜️ Data Compressor")
    col1, col2 = st.columns(2)
    with col1:
        comp_file = st.file_uploader("Upload PDF to compress", type=["pdf"])
        downscale = st.checkbox("Aggressive: Downscale all embedded images to 72 DPI")
        if comp_file and st.button("EXECUTE COMPRESSION"):
            with col2:
                with st.spinner("Running high-ratio compression..."):
                    input_path = save_uploaded_file(comp_file)
                    output_path = input_path.replace(".pdf", "_compressed.pdf")
                    try:
                        doc = fitz.open(input_path)
                        if downscale:
                            for page in doc:
                                for img in page.get_images(full=True):
                                    base_img = doc.extract_image(img[0])
                                    pil_img = Image.open(io.BytesIO(base_img["image"]))
                                    if pil_img.mode in ("RGBA", "P"): pil_img = pil_img.convert("RGB")
                                    w, h = pil_img.size
                                    pil_img = pil_img.resize((int(w*0.5), int(h*0.5)), Image.Resampling.LANCZOS)
                                    b = io.BytesIO()
                                    pil_img.save(b, format="JPEG", quality=60, optimize=True)
                                    page.insert_image(page.rect, stream=b.getvalue(), replace=True)
                                    
                        doc.save(output_path, garbage=4, deflate=True, clean=True)
                        doc.close()
                        st.success("COMPRESSION COMPLETE")
                        display_sizes(input_path, output_path)
                        with open(output_path, "rb") as f:
                            st.download_button("DOWNLOAD COMPRESSED", f, file_name="compressed.pdf")
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")

# ----------------- 6. DARK MODE INVERTER -----------------
with tabs[5]:
    st.header("🎨 Dark Mode Inverter")
    col1, col2 = st.columns(2)
    with col1:
        inv_file = st.file_uploader("Upload PDF", type=["pdf"], key="inv")
        dpi_setting = st.slider("Render Quality (DPI)", 72, 300, 150)
        if inv_file and st.button("INVERT SPECTRUM"):
            with col2:
                with st.spinner("Inverting visual spectrum..."):
                    input_path = save_uploaded_file(inv_file)
                    out_path = input_path.replace(".pdf", "_dark.pdf")
                    try:
                        doc = fitz.open(input_path)
                        out_pdf = fitz.open()
                        bar = st.progress(0)
                        for i, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=dpi_setting)
                            img = ImageOps.invert(Image.frombytes("RGB", [pix.width, pix.height], pix.samples))
                            b = io.BytesIO()
                            img.save(b, format="JPEG", quality=80)
                            out_pdf.insert_pdf(fitz.open("pdf", fitz.open("jpeg", b.getvalue()).convert_to_pdf()))
                            bar.progress((i+1)/len(doc))
                        out_pdf.save(out_path, deflate=True, garbage=4)
                        st.success("SPECTRUM INVERTED")
                        with open(out_path, "rb") as f:
                            st.download_button("DOWNLOAD DARK DOC", f, file_name="dark_mode.pdf")
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")

# ----------------- 7. AI REMBG -----------------
with tabs[6]:
    st.header("🖼️ Neural Background Eraser")
    col1, col2 = st.columns(2)
    with col1:
        bg_file = st.file_uploader("Upload PDF", type=["pdf"], key="bg")
        if bg_file and st.button("INITIALIZE NEURAL NET"):
            with col2:
                with st.spinner("Running U2Net Deep Learning model..."):
                    input_path = save_uploaded_file(bg_file)
                    out_path = input_path.replace(".pdf", "_nobg.pdf")
                    try:
                        doc = fitz.open(input_path)
                        out_pdf = fitz.open()
                        bar = st.progress(0)
                        for i, page in enumerate(doc):
                            pix = page.get_pixmap(dpi=150)
                            img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                            no_bg = remove(img)
                            bg = Image.new("RGB", no_bg.size, (10, 10, 10)) # dark background instead of white
                            bg.paste(no_bg, mask=no_bg.split()[3] if len(no_bg.split())==4 else None)
                            b = io.BytesIO()
                            bg.save(b, format="JPEG", quality=80)
                            out_pdf.insert_pdf(fitz.open("pdf", fitz.open("jpeg", b.getvalue()).convert_to_pdf()))
                            bar.progress((i+1)/len(doc))
                        out_pdf.save(out_path, deflate=True, garbage=4)
                        st.success("BACKGROUND ERADICATED")
                        with open(out_path, "rb") as f:
                            st.download_button("DOWNLOAD ARTIFACT", f, file_name="no_background.pdf")
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")

# ----------------- 8. SECURITY -----------------
with tabs[7]:
    st.header("🔐 Defense / Offense")
    col1, col2 = st.columns(2)
    with col1:
        sec_file = st.file_uploader("Upload Target Document", type=["pdf"])
        action = st.radio("Action Protocol", ["Encrypt (Defense)", "Decrypt (Offense)"])
        pwd = st.text_input("Enter Passphrase:", type="password")
        if sec_file and pwd and st.button("EXECUTE PROTOCOL"):
            with col2:
                with st.spinner("Processing cryptography..."):
                    input_path = save_uploaded_file(sec_file)
                    out_path = "secured.pdf"
                    try:
                        doc = fitz.open(input_path)
                        if action == "Decrypt (Offense)":
                            if doc.authenticate(pwd):
                                doc.save(out_path)
                                st.success("FIREWALL BREACHED. UNLOCKED.")
                            else:
                                st.error("ACCESS DENIED. Incorrect Passphrase.")
                                st.stop()
                        else:
                            doc.save(out_path, encryption=fitz.PDF_ENCRYPT_AES_256, user_pw=pwd, owner_pw=pwd)
                            st.success("ENCRYPTION ENABLED.")
                        doc.close()
                        with open(out_path, "rb") as f:
                            st.download_button("DOWNLOAD FILE", f, file_name="secured_doc.pdf")
                    except Exception as e:
                        st.error(f"SYSTEM FAULT: {e}")

st.markdown("<div class='custom-footer'>Made with <span>❤️</span> by <b>SONU VERMA</b></div>", unsafe_allow_html=True)
